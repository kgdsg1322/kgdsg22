#!/usr/bin/env python3
"""build_pptx.py — 설교용 PPTX 결정론 조립 (LLM이 코드를 매번 작성하지 않음)

구성: 표지(제목·본문구절) → 장면 슬라이드 ×N(이미지 풀블리드 + 상단 텍스트박스) → 마무리.
텍스트는 전부 텍스트박스 레이어 — 교사가 PPT에서 자유 수정 가능.
사용: python3 build_pptx.py <프로젝트폴더> <매핑표.md> --title <제목> [--verse <본문구절>]
출력: JSON. 슬라이드 수 = 장면 N + 2 검증 포함.
"""
import argparse, json, re, sys
from pathlib import Path


def mapping_rows(md: str):
    rows = []
    for line in md.splitlines():
        t = line.strip()
        if not t.startswith("|") or re.fullmatch(r"[|\s:\-]+", t):
            continue
        cells = [c.strip() for c in t.strip("|").split("|")]
        if len(cells) >= 5 and re.fullmatch(r"\d{2}_[^/\\\s]+\.jpg", cells[1]):
            rows.append({"file": cells[1], "scene": cells[2], "msg": cells[4]})
    return rows


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("project_dir")
    ap.add_argument("mapping")
    ap.add_argument("--title", required=True)
    ap.add_argument("--verse", default="")
    a = ap.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print(json.dumps({"status": "SKIP", "reason": "python-pptx 미설치 — pip3 install --user python-pptx 후 재실행"}, ensure_ascii=False))
        return 2

    proj = Path(a.project_dir).expanduser()
    rows = mapping_rows(Path(a.mapping).expanduser().read_text(encoding="utf-8", errors="replace"))
    if not rows:
        print(json.dumps({"status": "FAIL", "reason": "매핑표에서 장면 행을 찾지 못함"}, ensure_ascii=False))
        return 1
    missing = [r["file"] for r in rows if not (proj / "jpg" / r["file"]).exists()]
    if missing:
        print(json.dumps({"status": "FAIL", "reason": f"이미지 누락: {missing}"}, ensure_ascii=False))
        return 1

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def textbox(slide, x, y, w, h, text, size, bold=True, white=False, dark_bg=False, center=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        if dark_bg:
            tb.fill.solid()
            tb.fill.fore_color.rgb = RGBColor(28, 28, 28)
            tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        if center:
            p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        if white:
            r.font.color.rgb = RGBColor(255, 255, 255)
        return tb

    # 표지
    s = prs.slides.add_slide(blank)
    textbox(s, 1.0, 2.6, 11.3, 1.4, a.title, 54)
    if a.verse:
        textbox(s, 1.0, 4.2, 11.3, 1.0, a.verse, 30, bold=False)

    # 장면 슬라이드
    for r in rows:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(proj / "jpg" / r["file"]), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
        scene_name = re.sub(r"^\d{2}_|\.jpg$", "", r["file"])
        textbox(s, 0.4, 0.25, 12.5, 0.9, scene_name, 40, white=True, dark_bg=True)

    # 마무리
    s = prs.slides.add_slide(blank)
    textbox(s, 1.0, 3.0, 11.3, 1.4, "함께 기도해요", 44)

    out = proj / f"{a.title}_설교.pptx"
    prs.save(str(out))

    expected = len(rows) + 2
    actual = len(Presentation(str(out)).slides)
    status = "OK" if actual == expected and out.exists() else "FAIL"
    print(json.dumps({"status": status, "pptx": str(out),
                      "slides": actual, "expected": expected}, ensure_ascii=False))
    return 0 if status == "OK" else 1


if __name__ == "__main__":
    sys.exit(main())
